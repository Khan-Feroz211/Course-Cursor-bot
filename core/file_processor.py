import csv
import io
import logging
from pathlib import Path

import chardet
import fitz
import openpyxl
import pytesseract
from docx import Document
from PIL import Image, ImageEnhance, ImageOps
from pptx import Presentation

logger = logging.getLogger(__name__)


def extract_text(filepath: Path, file_type: str) -> list[dict]:
    pages: list[dict] = []
    source = filepath.name
    ext = file_type.lower().replace(".", "")
    try:
        if ext == "pdf":
            with fitz.open(filepath) as doc:
                for i, page in enumerate(doc, start=1):
                    txt = page.get_text("text").strip()
                    if txt:
                        pages.append({"text": txt, "page": i, "source": source})
        elif ext == "docx":
            doc = Document(str(filepath))
            paras = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
            for i in range(0, len(paras), 30):
                chunk = "\n".join(paras[i : i + 30]).strip()
                if chunk:
                    pages.append({"text": chunk, "page": (i // 30) + 1, "source": source})
        elif ext == "pptx":
            prs = Presentation(str(filepath))
            for slide_idx, slide in enumerate(prs.slides, start=1):
                parts = []
                for shp in slide.shapes:
                    if hasattr(shp, "text") and shp.text:
                        parts.append(shp.text.strip())
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes and notes.strip():
                        parts.append(f"Notes: {notes.strip()}")
                merged = "\n".join([x for x in parts if x]).strip()
                if merged:
                    pages.append({"text": merged, "page": slide_idx, "source": source})
        elif ext in {"xlsx", "xls"}:
            if ext == "xls":
                raise ValueError("Legacy .xls is not supported by current parser. Please save as .xlsx and upload again.")
            wb = openpyxl.load_workbook(str(filepath), data_only=True)
            for sheet in wb.worksheets:
                lines: list[str] = []
                for row in sheet.iter_rows(values_only=True):
                    vals = [str(v).strip() for v in row if v is not None and str(v).strip()]
                    if vals:
                        lines.append(" | ".join(vals))
                text = "\n".join(lines).strip()
                if text:
                    pages.append({"text": f"Sheet: {sheet.title}\n{text}", "page": len(pages) + 1, "source": source})
        elif ext == "csv":
            raw = filepath.read_bytes()
            enc = chardet.detect(raw).get("encoding") or "utf-8"
            content = raw.decode(enc, errors="replace")
            reader = csv.reader(io.StringIO(content))
            lines = []
            for row in reader:
                vals = [c.strip() for c in row if c and c.strip()]
                if vals:
                    lines.append(" | ".join(vals))
            text = "\n".join(lines).strip()
            if text:
                pages.append({"text": text, "page": 1, "source": source})
        elif ext == "txt":
            raw = filepath.read_bytes()
            enc = chardet.detect(raw).get("encoding") or "utf-8"
            text = raw.decode(enc, errors="replace").strip()
            if text:
                pages.append({"text": text, "page": 1, "source": source})
        elif ext in {"jpg", "jpeg", "png"}:
            img = Image.open(str(filepath))
            gray = ImageOps.grayscale(img)
            contrast = ImageEnhance.Contrast(gray).enhance(1.8)
            text = pytesseract.image_to_string(contrast).strip()
            if text:
                pages.append({"text": text, "page": 1, "source": source})
        else:
            logger.warning("File type not supported: %s", ext)
            raise ValueError("Unsupported file type")
    except Exception:
        logger.exception("File processing failure: %s", filepath.name)
        raise
    return pages


def chunk_pages(pages: list[dict], chunk_size: int = 400, overlap: int = 50) -> list[dict]:
    chunks: list[dict] = []
    if chunk_size <= overlap:
        raise ValueError("chunk_size must be greater than overlap")
    for page in pages:
        text = (page.get("text") or "").strip()
        if not text:
            continue
        words = text.split()
        start = 0
        chunk_index = 0
        while start < len(words):
            end = min(len(words), start + chunk_size)
            part = " ".join(words[start:end]).strip()
            if part:
                chunks.append(
                    {
                        "text": part,
                        "source_filename": page.get("source", "unknown"),
                        "page": int(page.get("page", 1)),
                        "chunk_index": chunk_index,
                    }
                )
                chunk_index += 1
            if end >= len(words):
                break
            start = max(0, end - overlap)
    return chunks
